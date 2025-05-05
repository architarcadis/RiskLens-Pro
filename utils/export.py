"""
Export module for RiskLens Pro
Provides functionality to export data, results, and visualizations to PDF and PPT
"""

import pandas as pd
import numpy as np
import io
import base64
import os
from typing import Dict, List, Tuple, Any, Union, Optional
from fpdf import FPDF
from io import BytesIO
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import tempfile
import plotly.graph_objects as go

# Constants
PRIMARY_COLOR = (0, 95, 158)  # RGB for #005f9e
SECONDARY_COLOR = (245, 166, 35)  # RGB for #f5a623
ACCENT_COLOR = (74, 74, 74)  # RGB for #4a4a4a


class PDFReport(FPDF):
    """
    Custom PDF class for RiskLens Pro reports
    """
    def __init__(self):
        super().__init__(orientation='P', unit='mm', format='A4')
        self.set_auto_page_break(auto=True, margin=15)
        self.set_margins(15, 15, 15)
        
        # Use default font (Helvetica)
        self.set_font('Helvetica', '', 10)
        
        # Add metadata
        self.set_title("RiskLens Pro Report")
        self.set_author("RiskLens Pro")
        self.set_creator("RiskLens Pro")
        
        # Add first page
        self.add_page()
    
    def header(self):
        """Define the header for all pages"""
        if self.page_no() > 1:  # Skip header on first page (cover page)
            # Set text color for header
            self.set_text_color(PRIMARY_COLOR[0], PRIMARY_COLOR[1], PRIMARY_COLOR[2])
            
            # Add logo-like element
            self.set_fill_color(PRIMARY_COLOR[0], PRIMARY_COLOR[1], PRIMARY_COLOR[2])
            self.rect(15, 10, 15, 10, style='F')
            
            # Add title
            self.set_font('Helvetica', 'B', 12)
            self.set_xy(35, 10)
            self.cell(0, 10, "RiskLens Pro Report", ln=1, align='L')
            
            # Add date
            self.set_font('Helvetica', '', 8)
            self.set_xy(35, 15)
            self.cell(0, 10, f"Generated on {datetime.now().strftime('%Y-%m-%d')}", ln=1, align='L')
            
            # Add horizontal line
            self.set_draw_color(200, 200, 200)
            self.line(15, 25, 195, 25)
            
            # Reset text color
            self.set_text_color(0, 0, 0)
            
            # Skip space after header
            self.ln(15)
    
    def footer(self):
        """Define the footer for all pages"""
        if self.page_no() > 1:  # Skip footer on first page (cover page)
            # Position at 1.5 cm from bottom
            self.set_y(-15)
            
            # Add horizontal line
            self.set_draw_color(200, 200, 200)
            self.line(15, self.get_y(), 195, self.get_y())
            
            # Set font
            self.set_font('Helvetica', '', 8)
            self.set_text_color(128, 128, 128)
            
            # Add page number and confidentiality notice
            self.cell(0, 10, f"Page {self.page_no()}", align='R')
            self.set_x(15)
            self.cell(0, 10, "Confidential | RiskLens Pro", align='L')
    
    def create_cover_page(self, title="Risk Analysis Report", subtitle=None, project_count=None):
        """Create a styled cover page"""
        # Set background color for top section
        self.set_fill_color(PRIMARY_COLOR[0], PRIMARY_COLOR[1], PRIMARY_COLOR[2])
        self.rect(0, 0, 210, 50, style='F')
        
        # Add logo placeholder
        self.set_fill_color(255, 255, 255)
        self.rect(15, 15, 30, 20, style='F')
        self.set_xy(15, 15)
        self.set_font('Helvetica', 'B', 18)
        self.set_text_color(PRIMARY_COLOR[0], PRIMARY_COLOR[1], PRIMARY_COLOR[2])
        self.cell(30, 20, "RL", ln=0, align='C')
        
        # Add title
        self.set_xy(50, 20)
        self.set_font('Helvetica', 'B', 24)
        self.set_text_color(255, 255, 255)
        self.cell(0, 10, title, ln=1, align='L')
        
        # Add subtitle if provided
        if subtitle:
            self.set_xy(50, 35)
            self.set_font('Helvetica', '', 14)
            self.cell(0, 10, subtitle, ln=1, align='L')
        
        # Add date
        self.set_xy(15, 70)
        self.set_font('Helvetica', '', 12)
        self.set_text_color(0, 0, 0)
        self.cell(0, 10, f"Generated on: {datetime.now().strftime('%Y-%m-%d')}", ln=1, align='L')
        
        # Add project count if provided
        if project_count is not None:
            self.set_xy(15, 85)
            self.cell(0, 10, f"Projects analyzed: {project_count}", ln=1, align='L')
        
        # Add description
        self.set_xy(15, 120)
        self.set_font('Helvetica', 'B', 14)
        self.set_text_color(PRIMARY_COLOR[0], PRIMARY_COLOR[1], PRIMARY_COLOR[2])
        self.cell(0, 10, "Executive Summary", ln=1, align='L')
        
        self.set_xy(15, 135)
        self.set_font('Helvetica', '', 11)
        self.set_text_color(0, 0, 0)
        description = (
            "This report provides a comprehensive analysis of project risk across your portfolio, "
            "powered by machine learning algorithms that identify patterns and risk factors. "
            "The insights and recommendations should be used to inform project management strategy "
            "and risk mitigation efforts."
        )
        self.multi_cell(0, 6, description, align='L')
        
        # Add confidentiality notice
        self.set_xy(15, 270)
        self.set_font('Helvetica', '', 8)
        self.set_text_color(128, 128, 128)
        self.cell(0, 10, "Confidential | For internal use only", ln=1, align='C')
    
    def add_section(self, title, spacing_before=10):
        """Add a section with styled heading"""
        # Add spacing
        self.ln(spacing_before)
        
        # Set text color for heading
        self.set_text_color(PRIMARY_COLOR[0], PRIMARY_COLOR[1], PRIMARY_COLOR[2])
        
        # Add section title
        self.set_font('Helvetica', 'B', 14)
        self.cell(0, 10, title, ln=1, align='L')
        
        # Add a line under the title
        self.set_draw_color(200, 200, 200)
        self.line(15, self.get_y(), 195, self.get_y())
        
        # Reset text color
        self.set_text_color(0, 0, 0)
        
        # Add a small space after the line
        self.ln(5)
    
    def add_subsection(self, title, spacing_before=5):
        """Add a subsection with styled heading"""
        # Add spacing
        self.ln(spacing_before)
        
        # Set text color for heading
        self.set_text_color(PRIMARY_COLOR[0], PRIMARY_COLOR[1], PRIMARY_COLOR[2])
        
        # Add subsection title
        self.set_font('Helvetica', 'B', 12)
        self.cell(0, 10, title, ln=1, align='L')
        
        # Reset text color
        self.set_text_color(0, 0, 0)
        
        # Add a small space after the title
        self.ln(2)
    
    def add_paragraph(self, text, spacing_before=2):
        """Add a paragraph of text"""
        # Add spacing
        self.ln(spacing_before)
        
        # Set font for paragraph
        self.set_font('Helvetica', '', 10)
        
        # Add paragraph text
        self.multi_cell(0, 5, text, align='L')
        
        # Add a small space after the paragraph
        self.ln(2)
    
    def add_metric(self, label, value, description=None, icon=None, spacing_before=5):
        """Add a metric with styled display"""
        # Add spacing
        self.ln(spacing_before)
        
        # Draw a box for the metric
        self.set_draw_color(200, 200, 200)
        self.set_fill_color(248, 248, 248)
        self.rect(15, self.get_y(), 180, 20, style='F')
        
        # Save current Y position
        current_y = self.get_y()
        
        # Add label
        self.set_xy(20, current_y + 2)
        self.set_font('Helvetica', '', 10)
        self.set_text_color(ACCENT_COLOR[0], ACCENT_COLOR[1], ACCENT_COLOR[2])
        self.cell(80, 8, label, ln=0, align='L')
        
        # Add value
        self.set_xy(100, current_y + 2)
        self.set_font('Helvetica', 'B', 12)
        self.set_text_color(PRIMARY_COLOR[0], PRIMARY_COLOR[1], PRIMARY_COLOR[2])
        self.cell(80, 8, str(value), ln=1, align='R')
        
        # Add description if provided
        if description:
            self.set_xy(20, current_y + 12)
            self.set_font('Helvetica', '', 8)
            self.set_text_color(128, 128, 128)
            self.multi_cell(170, 4, description, align='L')
        
        # Update Y position
        if description:
            self.set_y(current_y + 20)
        else:
            self.set_y(current_y + 20)
    
    def add_metrics_row(self, metrics, spacing_before=5):
        """Add a row of metrics side by side"""
        # Add spacing
        self.ln(spacing_before)
        
        # Calculate column width
        n_metrics = len(metrics)
        col_width = 180 / n_metrics
        
        # Save current position
        current_y = self.get_y()
        max_height = 0
        
        # Add each metric
        for i, metric in enumerate(metrics):
            # Draw a box for the metric
            x = 15 + i * col_width
            self.set_draw_color(200, 200, 200)
            self.set_fill_color(248, 248, 248)
            self.rect(x, current_y, col_width - 5, 30, style='F')
            
            # Add label
            self.set_xy(x + 5, current_y + 5)
            self.set_font('Helvetica', '', 9)
            self.set_text_color(ACCENT_COLOR[0], ACCENT_COLOR[1], ACCENT_COLOR[2])
            self.cell(col_width - 15, 6, metric['label'], ln=1, align='L')
            
            # Add value
            self.set_xy(x + 5, current_y + 15)
            self.set_font('Helvetica', 'B', 12)
            self.set_text_color(PRIMARY_COLOR[0], PRIMARY_COLOR[1], PRIMARY_COLOR[2])
            self.cell(col_width - 15, 8, str(metric['value']), ln=1, align='L')
        
        # Update Y position
        self.set_y(current_y + 35)
    
    def add_table(self, data, headers=None, spacing_before=5, col_widths=None):
        """Add a table with the given data"""
        # Add spacing
        self.ln(spacing_before)
        
        # Calculate column widths if not provided
        if not col_widths:
            n_cols = len(data[0]) if data else 0
            if n_cols > 0:
                col_widths = [180 / n_cols] * n_cols
        
        # Set font for table headers
        self.set_font('Helvetica', 'B', 10)
        self.set_fill_color(PRIMARY_COLOR[0], PRIMARY_COLOR[1], PRIMARY_COLOR[2])
        self.set_text_color(255, 255, 255)
        
        # Add headers if provided
        if headers:
            for i, header in enumerate(headers):
                self.cell(col_widths[i], 8, header, border=1, align='C', fill=True)
            self.ln()
        
        # Set font for table data
        self.set_font('Helvetica', '', 9)
        self.set_text_color(0, 0, 0)
        
        # Add table data with alternate row colors
        for j, row in enumerate(data):
            # Set background color for alternating rows
            if j % 2 == 0:
                self.set_fill_color(255, 255, 255)
            else:
                self.set_fill_color(240, 240, 240)
            
            # Check if we need to add a page break
            if self.get_y() + 8 > 270:  # Approaching bottom of page
                self.add_page()
                
                # Re-add headers after page break
                self.set_font('Helvetica', 'B', 10)
                self.set_fill_color(PRIMARY_COLOR[0], PRIMARY_COLOR[1], PRIMARY_COLOR[2])
                self.set_text_color(255, 255, 255)
                if headers:
                    for i, header in enumerate(headers):
                        self.cell(col_widths[i], 8, header, border=1, align='C', fill=True)
                    self.ln()
                
                # Reset font for table data
                self.set_font('Helvetica', '', 9)
                self.set_text_color(0, 0, 0)
                
                # Set background color again
                if j % 2 == 0:
                    self.set_fill_color(255, 255, 255)
                else:
                    self.set_fill_color(240, 240, 240)
            
            # Add each cell in the row
            for i, cell in enumerate(row):
                self.cell(col_widths[i], 8, str(cell), border=1, align='L', fill=True)
            self.ln()
    
    def add_plotly_figure(self, fig, width=180, height=100, spacing_before=5):
        """Add a plotly figure to the PDF"""
        # Add spacing
        self.ln(spacing_before)
        
        # Check if we need to add a page break
        if self.get_y() + height > 270:  # Approaching bottom of page
            self.add_page()
        
        # Save figure to a temporary file
        temp_dir = tempfile.gettempdir()
        img_path = os.path.join(temp_dir, 'temp_plot.png')
        fig.write_image(img_path, width=width*3, height=height*3, scale=3)
        
        # Add image to PDF
        self.image(img_path, x=15, y=None, w=width)
        
        # Update Y position
        self.ln(height + 5)
        
        # Clean up temporary file
        try:
            os.remove(img_path)
        except:
            pass


def create_pdf_report(project_data, model_results, visualizations, risk_data=None):
    """
    Create a PDF report with project data, model results, and visualizations
    Args:
        project_data: DataFrame with project data
        model_results: Dictionary of model results
        visualizations: Dictionary of visualization figures
        risk_data: DataFrame with risk data (optional)
    Returns:
        BytesIO: PDF report as BytesIO object
    """
    # Create PDF object
    pdf = PDFReport()
    
    # Create cover page
    pdf.create_cover_page(
        title="Project Risk Analysis Report",
        subtitle="Powered by RiskLens Pro",
        project_count=len(project_data) if project_data is not None else None
    )
    
    # Add Executive Summary
    pdf.add_page()
    pdf.add_section("Executive Summary")
    
    # Add narrative summary if available
    if 'narrative_summary' in visualizations:
        pdf.add_paragraph(visualizations['narrative_summary'])
    else:
        pdf.add_paragraph(
            "This report presents the results of a machine learning-based risk analysis "
            "of your project portfolio. It identifies potential risk factors and provides "
            "recommendations for risk mitigation strategies."
        )
    
    # Add key metrics
    if project_data is not None and 'ProjectDerailmentRisk' in project_data.columns:
        pdf.add_subsection("Key Metrics")
        
        high_risk = project_data['ProjectDerailmentRisk'].sum()
        total_projects = len(project_data)
        high_risk_rate = high_risk / total_projects * 100 if total_projects > 0 else 0
        
        metrics = [
            {'label': 'Total Projects', 'value': total_projects},
            {'label': 'High Risk Projects', 'value': high_risk},
            {'label': 'High Risk Rate', 'value': f"{high_risk_rate:.1f}%"}
        ]
        
        pdf.add_metrics_row(metrics)
    
    # Add model performance metrics if available
    if model_results is not None and 'metrics' in model_results and model_results['best_model'] is not None:
        best_model = model_results['best_model']
        model_metrics = model_results['metrics'].get(best_model, {})
        
        if model_metrics:
            pdf.add_subsection(f"Model Performance Metrics ({best_model})")
            
            metrics = [
                {'label': 'Accuracy', 'value': f"{model_metrics.get('accuracy', 0):.2f}"},
                {'label': 'Precision', 'value': f"{model_metrics.get('precision', 0):.2f}"},
                {'label': 'Recall', 'value': f"{model_metrics.get('recall', 0):.2f}"},
                {'label': 'ROC AUC', 'value': f"{model_metrics.get('roc_auc', 0):.2f}"}
            ]
            
            pdf.add_metrics_row(metrics)
    
    # Add feature importance if available
    if 'feature_importance_fig' in visualizations:
        pdf.add_page()
        pdf.add_section("Feature Importance Analysis")
        pdf.add_paragraph(
            "The following chart shows the most important factors in predicting project risk. "
            "These factors have the strongest influence on the model's predictions."
        )
        pdf.add_plotly_figure(visualizations['feature_importance_fig'])
    
    # Add Project Risk Analysis
    if project_data is not None and 'RiskProbability' in project_data.columns:
        pdf.add_page()
        pdf.add_section("Project Risk Analysis")
        
        pdf.add_subsection("Top High-Risk Projects")
        # Create table of top 10 high-risk projects
        top_risk_projects = project_data.sort_values('RiskProbability', ascending=False).head(10)
        
        headers = ["Project ID", "Project Name", "Risk Probability", "Project Type"]
        data = []
        for _, row in top_risk_projects.iterrows():
            project_id = row.get('ProjectID', 'N/A')
            project_name = row.get('ProjectName', 'N/A')
            risk_prob = f"{row.get('RiskProbability', 0) * 100:.1f}%"
            project_type = row.get('ProjectType', 'N/A')
            
            data.append([project_id, project_name, risk_prob, project_type])
        
        col_widths = [30, 70, 30, 50]
        pdf.add_table(data, headers, col_widths=col_widths)
    
    # Add risk distribution visualization if available
    if 'risk_distribution_fig' in visualizations:
        pdf.add_page()
        pdf.add_section("Risk Distribution")
        pdf.add_paragraph(
            "The following chart shows the distribution of risk probabilities across your project portfolio. "
            "This helps to understand the overall risk landscape."
        )
        pdf.add_plotly_figure(visualizations['risk_distribution_fig'])
    
    # Add Risk Register Analysis if available
    if risk_data is not None and len(risk_data) > 0:
        pdf.add_page()
        pdf.add_section("Risk Register Analysis")
        pdf.add_paragraph(
            "The following section provides an analysis of your risk register data, "
            "showing the distribution of risks by type, impact, and probability."
        )
        
        # Add Risk Register metrics
        risk_count = len(risk_data)
        open_risks = len(risk_data[risk_data['Status'] == 'Open']) if 'Status' in risk_data.columns else 'N/A'
        
        metrics = [
            {'label': 'Total Risks', 'value': risk_count},
            {'label': 'Open Risks', 'value': open_risks},
        ]
        
        if 'RiskType' in risk_data.columns:
            most_common_type = risk_data['RiskType'].value_counts().index[0] if len(risk_data) > 0 else 'N/A'
            metrics.append({'label': 'Most Common Risk Type', 'value': most_common_type})
        
        pdf.add_metrics_row(metrics)
        
        # Add Risk Register heatmap if available
        if 'risk_heatmap_fig' in visualizations:
            pdf.add_plotly_figure(visualizations['risk_heatmap_fig'])
    
    # Add recommendations
    pdf.add_page()
    pdf.add_section("Recommendations")
    pdf.add_paragraph(
        "Based on the analysis, the following recommendations are provided to mitigate project risks:"
    )
    
    # Generate recommendations based on available data
    recommendations = [
        "Focus on high-risk projects identified in the report, particularly those with the highest risk probabilities.",
        "Address the key risk factors identified in the feature importance analysis.",
        "Implement regular risk assessment reviews for projects with changing risk profiles.",
        "Develop specific mitigation strategies for the most common risk types in your portfolio.",
        "Consider reallocating resources to support projects at highest risk of derailment."
    ]
    
    # Add recommendations as bullet points
    for i, recommendation in enumerate(recommendations, 1):
        pdf.add_paragraph(f"  {i}. {recommendation}")
    
    # Add conclusion
    pdf.add_section("Conclusion", spacing_before=15)
    pdf.add_paragraph(
        "This report provides a comprehensive analysis of your project portfolio's risk landscape. "
        "By leveraging machine learning techniques, we have identified key risk factors and projects "
        "that require attention. The recommendations provided should be integrated into your "
        "project management and risk mitigation strategies to improve overall project success rates."
    )
    
    # Create BytesIO object to save the PDF to
    pdf_buffer = BytesIO()
    pdf.output(dest='F', name=pdf_buffer)
    pdf_buffer.seek(0)
    
    return pdf_buffer


def create_ppt_report(project_data, model_results, visualizations, risk_data=None):
    """
    Create a PowerPoint presentation with project data, model results, and visualizations
    Args:
        project_data: DataFrame with project data
        model_results: Dictionary of model results
        visualizations: Dictionary of visualization figures
        risk_data: DataFrame with risk data (optional)
    Returns:
        BytesIO: PowerPoint presentation as BytesIO object
    """
    # Create presentation
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]  # Title Slide
    title_content_layout = prs.slide_layouts[1]  # Title and Content
    section_layout = prs.slide_layouts[2]  # Section Header
    two_content_layout = prs.slide_layouts[3]  # Two Content
    
    # Add title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Project Risk Analysis Report"
    subtitle.text = f"RiskLens Pro | {datetime.now().strftime('%Y-%m-%d')}"
    
    # Add executive summary slide
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Executive Summary"
    
    # Add executive summary text
    if 'narrative_summary' in visualizations:
        content.text = visualizations['narrative_summary']
    else:
        content.text = (
            "This presentation summarizes the results of a machine learning-based risk analysis "
            "of your project portfolio. It identifies potential risk factors and provides "
            "recommendations for risk mitigation strategies."
        )
    
    # Add key metrics slide
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "Key Metrics"
    
    # Create a table for metrics
    if project_data is not None and 'ProjectDerailmentRisk' in project_data.columns:
        high_risk = project_data['ProjectDerailmentRisk'].sum()
        total_projects = len(project_data)
        high_risk_rate = high_risk / total_projects * 100 if total_projects > 0 else 0
        
        rows, cols = 2, 3
        left = Inches(1.0)
        top = Inches(2.0)
        width = Inches(8.0)
        height = Inches(1.5)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set headers
        table.cell(0, 0).text = "Total Projects"
        table.cell(0, 1).text = "High Risk Projects"
        table.cell(0, 2).text = "High Risk Rate"
        
        # Set values
        table.cell(1, 0).text = str(total_projects)
        table.cell(1, 1).text = str(high_risk)
        table.cell(1, 2).text = f"{high_risk_rate:.1f}%"
        
        # Style the table
        for i in range(2):
            for j in range(3):
                cell = table.cell(i, j)
                if i == 0:  # Header row
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(PRIMARY_COLOR[0], PRIMARY_COLOR[1], PRIMARY_COLOR[2])
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                        paragraph.font.bold = True
                else:  # Data row
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(14)
    
    # Add model performance slide if available
    if model_results is not None and 'metrics' in model_results and model_results['best_model'] is not None:
        slide = prs.slides.add_slide(title_content_layout)
        title = slide.shapes.title
        title.text = "Model Performance"
        
        best_model = model_results['best_model']
        model_metrics = model_results['metrics'].get(best_model, {})
        
        if model_metrics:
            # Create a table for metrics
            rows, cols = 2, 4
            left = Inches(1.0)
            top = Inches(2.0)
            width = Inches(8.0)
            height = Inches(1.5)
            
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Set headers
            table.cell(0, 0).text = "Accuracy"
            table.cell(0, 1).text = "Precision"
            table.cell(0, 2).text = "Recall"
            table.cell(0, 3).text = "ROC AUC"
            
            # Set values
            table.cell(1, 0).text = f"{model_metrics.get('accuracy', 0):.2f}"
            table.cell(1, 1).text = f"{model_metrics.get('precision', 0):.2f}"
            table.cell(1, 2).text = f"{model_metrics.get('recall', 0):.2f}"
            table.cell(1, 3).text = f"{model_metrics.get('roc_auc', 0):.2f}"
            
            # Style the table
            for i in range(2):
                for j in range(4):
                    cell = table.cell(i, j)
                    if i == 0:  # Header row
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(PRIMARY_COLOR[0], PRIMARY_COLOR[1], PRIMARY_COLOR[2])
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.color.rgb = RGBColor(255, 255, 255)
                            paragraph.font.bold = True
                    else:  # Data row
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.bold = True
                            paragraph.font.size = Pt(14)
            
            # Add model info
            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(8.0), Inches(1.0))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = f"Best Model: {best_model}"
            p.font.bold = True
    
    # Add feature importance slide if available
    if 'feature_importance_fig' in visualizations:
        slide = prs.slides.add_slide(title_content_layout)
        title = slide.shapes.title
        title.text = "Feature Importance Analysis"
        
        # Save figure to temp file
        temp_dir = tempfile.gettempdir()
        img_path = os.path.join(temp_dir, 'feature_importance.png')
        visualizations['feature_importance_fig'].write_image(img_path, width=1000, height=600)
        
        # Add image to slide
        content = slide.placeholders[1]
        content.insert_picture(img_path)
        
        # Clean up temporary file
        try:
            os.remove(img_path)
        except:
            pass
    
    # Add high-risk projects slide
    if project_data is not None and 'RiskProbability' in project_data.columns:
        slide = prs.slides.add_slide(title_content_layout)
        title = slide.shapes.title
        title.text = "Top High-Risk Projects"
        
        # Create table of top 5 high-risk projects
        top_risk_projects = project_data.sort_values('RiskProbability', ascending=False).head(5)
        
        rows, cols = len(top_risk_projects) + 1, 4
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(9.0)
        height = Inches(3.0)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set headers
        table.cell(0, 0).text = "Project ID"
        table.cell(0, 1).text = "Project Name"
        table.cell(0, 2).text = "Risk Probability"
        table.cell(0, 3).text = "Project Type"
        
        # Set values
        for i, (_, row) in enumerate(top_risk_projects.iterrows(), 1):
            table.cell(i, 0).text = str(row.get('ProjectID', 'N/A'))
            table.cell(i, 1).text = str(row.get('ProjectName', 'N/A'))
            table.cell(i, 2).text = f"{row.get('RiskProbability', 0) * 100:.1f}%"
            table.cell(i, 3).text = str(row.get('ProjectType', 'N/A'))
        
        # Style the table
        for i in range(rows):
            for j in range(cols):
                cell = table.cell(i, j)
                if i == 0:  # Header row
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(PRIMARY_COLOR[0], PRIMARY_COLOR[1], PRIMARY_COLOR[2])
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.color.rgb = RGBColor(255, 255, 255)
                        paragraph.font.bold = True
    
    # Add risk distribution slide if available
    if 'risk_distribution_fig' in visualizations:
        slide = prs.slides.add_slide(title_content_layout)
        title = slide.shapes.title
        title.text = "Risk Distribution"
        
        # Save figure to temp file
        temp_dir = tempfile.gettempdir()
        img_path = os.path.join(temp_dir, 'risk_distribution.png')
        visualizations['risk_distribution_fig'].write_image(img_path, width=1000, height=600)
        
        # Add image to slide
        content = slide.placeholders[1]
        content.insert_picture(img_path)
        
        # Clean up temporary file
        try:
            os.remove(img_path)
        except:
            pass
    
    # Add recommendations slide
    slide = prs.slides.add_slide(title_content_layout)
    title = slide.shapes.title
    title.text = "Recommendations"
    
    # Generate recommendations based on available data
    recommendations = [
        "Focus on high-risk projects identified in the report, particularly those with the highest risk probabilities.",
        "Address the key risk factors identified in the feature importance analysis.",
        "Implement regular risk assessment reviews for projects with changing risk profiles.",
        "Develop specific mitigation strategies for the most common risk types in your portfolio.",
        "Consider reallocating resources to support projects at highest risk of derailment."
    ]
    
    # Add recommendations to slide
    content = slide.placeholders[1]
    tf = content.text_frame
    
    for i, recommendation in enumerate(recommendations, 1):
        p = tf.add_paragraph()
        p.text = f"{i}. {recommendation}"
        p.level = 0
    
    # Create BytesIO object to save the presentation to
    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    
    return ppt_buffer


def get_download_link(file_bytes, file_name, file_format="PDF"):
    """
    Create a download link for a file
    Args:
        file_bytes: BytesIO object containing the file
        file_name: Name of the file
        file_format: Format of the file (PDF or PPTX)
    Returns:
        str: HTML download link
    """
    # Determine MIME type based on format
    if file_format.upper() == "PDF":
        mime_type = "application/pdf"
    elif file_format.upper() == "PPTX":
        mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    else:
        mime_type = "application/octet-stream"
    
    # Encode file as base64
    b64 = base64.b64encode(file_bytes.read()).decode()
    
    # Create download link
    href = f'<a href="data:{mime_type};base64,{b64}" download="{file_name}" target="_blank"><button style="background-color: #005f9e; color: white; padding: 8px 15px; border: none; border-radius: 4px; cursor: pointer; font-weight: bold; width: 100%;">Download {file_format}</button></a>'
    
    return href
